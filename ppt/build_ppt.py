# -*- coding: utf-8 -*-
"""口試簡報生成器 — SACF 情感感知跨模態融合框架與多代理模擬"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")

# ---- 色彩主題 ----
NAVY   = RGBColor(0x16, 0x2A, 0x4D)   # 主深藍
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)   # 強調藍
TEAL   = RGBColor(0x12, 0x84, 0x8D)   # 輔色青
RED    = RGBColor(0xD6, 0x3A, 0x3A)   # 重點紅（本研究/結果）
GOLD   = RGBColor(0xE0, 0x9F, 0x3E)
LIGHT  = RGBColor(0xF2, 0xF5, 0xFA)   # 淺底
GREY   = RGBColor(0x55, 0x5B, 0x66)
DARK   = RGBColor(0x22, 0x2A, 0x38)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARDBG = RGBColor(0xEA, 0xF1, 0xF9)

FONT = "PingFang TC"
FONT_NUM = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_font(run, size=None, bold=None, color=None, name=FONT):
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True):
    """lines: list of dicts: {t, size, bold, color, name, space_after, bullet, level, space_before}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if "space_after" in ln: p.space_after = Pt(ln["space_after"])
        if "space_before" in ln: p.space_before = Pt(ln["space_before"])
        if "line_spacing" in ln: p.line_spacing = ln["line_spacing"]
        segs = ln["t"] if isinstance(ln["t"], list) else [ln]
        for seg in segs:
            r = p.add_run(); r.text = seg.get("t", "")
            _set_font(r, seg.get("size", ln.get("size", 18)),
                      seg.get("bold", ln.get("bold", False)),
                      seg.get("color", ln.get("color", DARK)),
                      seg.get("name", ln.get("name", FONT)))
    return tb


def pic_fit(slide, path, x, y, max_w, max_h, center=True, border=True):
    im = Image.open(path); iw, ih = im.size
    ar = iw / ih
    w = max_w; h = Emu(int(w / ar))
    if h > max_h:
        h = max_h; w = Emu(int(h * ar))
    px = x + (max_w - w) // 2 if center else x
    py = y + (max_h - h) // 2 if center else y
    pic = slide.shapes.add_picture(path, px, py, w, h)
    if border:
        pic.line.color.rgb = RGBColor(0xC9, 0xD3, 0xE0); pic.line.width = Pt(0.75)
    return pic


# ---- 版面元件 ----
def content_header(slide, chapter, title):
    rect(slide, 0, 0, SW, Inches(1.0), NAVY)
    rect(slide, 0, Inches(1.0), SW, Pt(4), GOLD)
    if chapter:
        txt(slide, Inches(0.5), Inches(0.10), Inches(4), Inches(0.3),
            [{"t": chapter, "size": 13, "bold": True, "color": GOLD}])
    txt(slide, Inches(0.5), Inches(0.34), Inches(12.4), Inches(0.62),
        [{"t": title, "size": 25, "bold": True, "color": WHITE}],
        anchor=MSO_ANCHOR.MIDDLE)
    # 頁碼於 footer 由 add_footer 處理


def add_footer(slide, idx):
    txt(slide, Inches(11.8), Inches(7.05), Inches(1.4), Inches(0.35),
        [{"t": "SACF · 李昇峰", "size": 9, "color": GREY}], align=PP_ALIGN.RIGHT)
    txt(slide, Inches(0.4), Inches(7.05), Inches(1.0), Inches(0.35),
        [{"t": str(idx), "size": 10, "color": GREY, "bold": True}])


def bullet(t, size=17, color=DARK, bold=False, sa=10, sub=False):
    mark = "▪ " if not sub else "    – "
    return {"t": [{"t": mark, "size": size, "color": BLUE if not sub else TEAL, "bold": True},
                  {"t": t, "size": size, "color": color, "bold": bold}],
            "space_after": sa, "line_spacing": 1.12}


def card(slide, x, y, w, h, head, body, headcolor=BLUE, bg=CARDBG):
    rect(slide, x, y, w, h, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, x, y, Pt(6), h, headcolor)
    lines = [{"t": head, "size": 15, "bold": True, "color": headcolor, "space_after": 5}]
    for b in body:
        lines.append({"t": b, "size": 12.5, "color": DARK, "space_after": 3, "line_spacing": 1.08})
    txt(slide, x + Inches(0.18), y + Inches(0.10), w - Inches(0.3), h - Inches(0.2), lines)


def metric_chip(slide, x, y, w, value, label, vcolor=RED):
    h = Inches(1.05)
    rect(slide, x, y, w, h, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, x, y, w, Pt(5), vcolor)
    txt(slide, x, y + Inches(0.12), w, Inches(0.55),
        [{"t": value, "size": 27, "bold": True, "color": vcolor, "name": FONT_NUM}],
        align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(0.66), w, Inches(0.35),
        [{"t": label, "size": 12, "color": GREY}], align=PP_ALIGN.CENTER)


P = 0  # 全域頁碼


def page(slide):
    global P
    P += 1
    add_footer(slide, P)


# =====================================================================
# Slide 1 — 封面
# =====================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.55), SW, Inches(2.35), RGBColor(0x1B, 0x35, 0x60))
rect(s, 0, Inches(2.55), SW, Pt(4), GOLD)
rect(s, 0, Inches(4.86), SW, Pt(4), GOLD)
# logos
for lp, lx in [("image1.png", Inches(5.95)),]:
    p = os.path.join(ASSETS, lp)
    if os.path.exists(p):
        pic_fit(s, p, lx, Inches(0.45), Inches(1.45), Inches(1.45), border=False)
txt(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.5),
    [{"t": "淡江大學  統計與資料科學碩士班  碩士論文", "size": 17, "color": GOLD, "bold": True}],
    align=PP_ALIGN.CENTER)
txt(s, Inches(0.6), Inches(2.62), Inches(12.1), Inches(1.5),
    [{"t": "SACF：情感感知跨模態融合框架與極性增強注意力", "size": 29, "bold": True, "color": WHITE,
      "space_after": 6, "align": PP_ALIGN.CENTER},
     {"t": "模型建構及多代理模擬應用研究", "size": 29, "bold": True, "color": WHITE,
      "align": PP_ALIGN.CENTER}])
txt(s, Inches(0.6), Inches(4.12), Inches(12.1), Inches(0.6),
    [{"t": "A Sentiment-Aware Cross-modal Fusion Framework with Polarity-Enhanced",
      "size": 13, "color": RGBColor(0xC6, 0xD3, 0xE6), "align": PP_ALIGN.CENTER, "space_after": 2},
     {"t": "Attention and Its Application to Multi-Agent Simulation",
      "size": 13, "color": RGBColor(0xC6, 0xD3, 0xE6), "align": PP_ALIGN.CENTER}])
txt(s, Inches(0.6), Inches(5.25), Inches(12.1), Inches(1.4),
    [{"t": "研究生：李昇峰　撰", "size": 19, "bold": True, "color": WHITE,
      "align": PP_ALIGN.CENTER, "space_after": 8},
     {"t": "指導教授：高君豪　教授", "size": 17, "color": RGBColor(0xD9, 0xE2, 0xF0),
      "align": PP_ALIGN.CENTER, "space_after": 8},
     {"t": "中華民國 115 年 6 月", "size": 14, "color": GOLD, "align": PP_ALIGN.CENTER}])
page(s)

# =====================================================================
# Slide 2 — 簡報大綱
# =====================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "OUTLINE", "簡報大綱")
items = [
    ("01", "研究背景與動機", "情感建模缺口 · 雙核心目標", BLUE),
    ("02", "文獻探討", "智能體理論 · 多模態情感分析演進", TEAL),
    ("03", "SACF 模型方法", "架構 · 損失函數 · 訓練與推斷", NAVY),
    ("04", "實驗結果", "CMU-MOSI · 跨資料集泛化 · 參數效率", RED),
    ("05", "情緒感知多代理模擬系統", "認知循環 · 情緒子系統 · 系統展示", TEAL),
    ("06", "結論與貢獻", "研究成果 · 三大貢獻 · 未來展望", GOLD),
]
x0 = Inches(0.7); y0 = Inches(1.45)
cw = Inches(5.95); ch = Inches(1.55); gx = Inches(0.35); gy = Inches(0.3)
for i, (no, t, d, c) in enumerate(items):
    r, col = divmod(i, 2)
    x = x0 + col * (cw + gx); y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, Inches(1.1), ch, c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, y, Inches(1.1), ch, [{"t": no, "size": 30, "bold": True, "color": WHITE}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(1.3), y + Inches(0.22), cw - Inches(1.45), ch - Inches(0.3),
        [{"t": t, "size": 19, "bold": True, "color": DARK, "space_after": 5},
         {"t": d, "size": 13, "color": GREY}])
page(s)


# =====================================================================
# 章節分隔頁
# =====================================================================
def divider(no, title, sub, points):
    s = add_slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 0, Inches(4.6), SH, RGBColor(0x10, 0x20, 0x3C))
    txt(s, Inches(0.55), Inches(2.5), Inches(3.6), Inches(2),
        [{"t": "CHAPTER", "size": 15, "color": GOLD, "bold": True, "space_after": 4},
         {"t": no, "size": 110, "bold": True, "color": WHITE}])
    txt(s, Inches(5.05), Inches(2.55), Inches(7.7), Inches(1.2),
        [{"t": title, "size": 36, "bold": True, "color": WHITE, "space_after": 6},
         {"t": sub, "size": 16, "color": GOLD}])
    ly = Inches(4.15)
    for p in points:
        txt(s, Inches(5.05), ly, Inches(7.7), Inches(0.45),
            [{"t": [{"t": "▸ ", "size": 16, "color": GOLD, "bold": True},
                    {"t": p, "size": 16, "color": RGBColor(0xD9, 0xE2, 0xF0)}]}])
        ly += Inches(0.5)
    page(s)
    return s


# =====================================================================
# CHAPTER 1
# =====================================================================
divider("1", "研究背景與動機", "Background & Motivation",
        ["生成式代理的情感建模缺口", "雙核心研究目標", "研究貢獻概覽"])

# 1.1 背景與動機
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第一章 · 研究背景", "為什麼需要「會感受情緒」的 AI 代理？")
txt(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.8),
    [{"t": [{"t": "Park 等人（2023）", "size": 16, "bold": True, "color": BLUE},
            {"t": "的生成式代理首次驗證 LLM 能驅動具記憶、反思、規劃能力的虛擬居民，模擬逼真社會行為；", "size": 16, "color": DARK},
            {"t": "然而其情感狀態僅靠文字語義推理，缺乏量化表示與動態更新機制。", "size": 16, "bold": True, "color": RED}],
      "line_spacing": 1.15}])
cards = [
    ("① 情感建模缺口", ["行為決策完全依賴文字語義推理", "缺乏情感的量化表示與動態更新", "互動真實性存在根本缺陷"], RED),
    ("② 社交真實性需求", ["情緒是人際互動的核心驅動力", "情緒影響溝通風格與社交距離", "缺情緒回應 → 機械化文字交換"], BLUE),
    ("③ 技術突破可行", ["DeBERTa 解耦注意力捕捉情感語義", "BiLSTM 編碼音訊/視覺時序", "跨模態注意力融合已成熟"], TEAL),
    ("④ 智能問卷創新", ["傳統問卷受記憶偏差與社會期望干擾", "情緒感知代理可一致且可追溯", "支援反覆採樣與假設驗證"], GOLD),
]
cx = Inches(0.6); cw = Inches(2.95); cy = Inches(2.35); ch = Inches(2.7); gx = Inches(0.18)
for i, (h, b, c) in enumerate(cards):
    card(s, cx + i * (cw + gx), cy, cw, ch, h, b, headcolor=c)
txt(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.3),
    [{"t": "研究切入點", "size": 14, "bold": True, "color": NAVY, "space_after": 3},
     {"t": "整合「多模態情感分析」與「生成式代理模擬」兩大技術線索 —— 以 SACF 情感評分引擎為 9 位 AI 居民提供可量化的情緒感知基礎，並利用情緒記憶自動填寫問卷。",
      "size": 15, "color": DARK, "line_spacing": 1.15}])
page(s)

# 1.2 研究目的
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第一章 · 研究目的", "雙核心目標：情緒感知模擬 × 智能問卷調查")
# 左右兩大目標
rect(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(0.6), BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(0.6),
    [{"t": "目標 A｜建構 SACF 情感評分模型", "size": 16, "bold": True, "color": WHITE}],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(0.6), TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(0.6),
    [{"t": "目標 B｜情緒感知多代理模擬系統", "size": 16, "bold": True, "color": WHITE}],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.7), Inches(2.05), Inches(5.9), Inches(4.5),
    [bullet("以 CMU-MOSI 為基礎開發跨模態融合框架", 15),
     bullet("整合 DeBERTa-v3-large + BiLSTM 編碼器", 14, sub=True),
     bullet("PEA 極性增強注意力 + 情感感知融合", 14, sub=True),
     bullet("多工損失 + 兩階段訓練 + 參數空間集成", 14, sub=True),
     bullet("七分類情感強度 Acc-7 達 53.06%", 15, color=RED, bold=True),
     bullet("情緒驅動的智能體互動機制", 15),
     bullet("以情感分數動態調整對話與行為規劃", 14, sub=True)])
txt(s, Inches(6.95), Inches(2.05), Inches(5.9), Inches(4.5),
    [bullet("9 位 AI 居民人格塑造與情感建模", 15),
     bullet("七類情緒狀態轉換機制（正/負向）", 14, sub=True),
     bullet("三維記憶管理（時效/重要/相關）", 15),
     bullet("衰減因子 0.995 · LLM 1–10 分重要性", 14, sub=True),
     bullet("虛擬環境與社交網絡視覺化", 15),
     bullet("智能問卷系統（單選/多選/評分/文字）", 15),
     bullet("多格式匯出 CSV / JSON / Excel", 14, sub=True)])
page(s)


# =====================================================================
# CHAPTER 2
# =====================================================================
divider("2", "文獻探討", "Literature Review",
        ["智能體理論基礎演進", "多模態情感分析三條方法論路線", "與最強基線 MGT 的差異化設計"])

# 2.1 學術脈絡
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第二章 · 文獻脈絡", "四條交匯的學術脈絡")
flow = [("智能體理論", "感知-推理-行動\n→ 情感三元模型", BLUE),
        ("情感計算", "Ekman 離散情感\nRussell 環形連續維度", TEAL),
        ("多智能體系統", "社會互動建模\n計算社會科學", GOLD),
        ("LLM 智能體", "生成式代理\n人格化與角色扮演", NAVY)]
fx = Inches(0.7); fw = Inches(2.75); fy = Inches(1.55); fh = Inches(1.7); gap = Inches(0.35)
for i, (h, d, c) in enumerate(flow):
    x = fx + i * (fw + gap)
    rect(s, x, fy, fw, fh, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, fy, fw, Inches(0.5), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, fy, fw, Inches(0.5), [{"t": h, "size": 15, "bold": True, "color": WHITE}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, fy + Inches(0.6), fw, Inches(1.0),
        [{"t": d, "size": 13, "color": DARK, "line_spacing": 1.1}], align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, x + fw, fy + Inches(0.5), gap, Inches(0.7),
            [{"t": "+", "size": 22, "bold": True, "color": GREY}], align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(3.55), Inches(11.9), Inches(0.6),
    [{"t": "共同指向同一個尚未被系統性解決的問題：", "size": 15, "color": GREY},
     ], align=PP_ALIGN.CENTER)
rect(s, Inches(1.6), Inches(4.05), Inches(10.1), Inches(0.85), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.6), Inches(4.05), Inches(10.1), Inches(0.85),
    [{"t": "如何在生成式多代理模擬環境中，以「可量化的情感感知能力」驅動更真實的社交互動？",
      "size": 17, "bold": True, "color": NAVY}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.6),
    [bullet("Russell 情感環形模型 → 啟發 SACF 採七點量表（−3~+3）連續強度，捕捉細微差異", 14.5),
     bullet("MSA 三條路線：自監督表示學習 · 對比學習重塑特徵空間 · 表示空間分解提純情感訊號", 14.5),
     bullet("典範轉移：文字模態成為「語義主幹」，音訊/視覺降為補充注入的「情感線索」", 14.5)])
page(s)

# 2.2 與 MGT 差異化
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第二章 · 研究定位", "SACF 與最強基線 MGT 的三點差異化設計")
# 表頭
hx = [Inches(0.6), Inches(3.4), Inches(7.3), Inches(10.6)]
hw = [Inches(2.8), Inches(3.9), Inches(3.3), Inches(2.5)]
heads = ["設計維度", "MGT（對標基線）", "本研究 SACF", ""]
ty = Inches(1.4)
rows = [
    ("骨幹模型", "生成式 T5", "判別式 DeBERTa-v3-large\n解耦注意力精準捕捉情感詞元"),
    ("查詢向量", "整體語言表示作查詢", "PEA 情感顯著詞元構建\n情感感知查詢「定點對準」"),
    ("訓練目標", "邊際損失", "序數 EMD + 軟序數標籤\n+ R-Drop + 跨模態 InfoNCE"),
]
# header
rect(s, Inches(0.6), ty, Inches(12.1), Inches(0.55), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
for i, h in enumerate(["設計維度", "MGT（對標最強基線）", "本研究 SACF 差異化設計"]):
    xx = [Inches(0.6), Inches(3.5), Inches(7.6)][i]
    ww = [Inches(2.9), Inches(4.1), Inches(5.1)][i]
    txt(s, xx, ty, ww, Inches(0.55), [{"t": h, "size": 15, "bold": True, "color": WHITE}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
ry = ty + Inches(0.65)
for j, (a, b, c) in enumerate(rows):
    rh = Inches(0.95)
    bg = WHITE if j % 2 == 0 else CARDBG
    rect(s, Inches(0.6), ry, Inches(12.1), rh, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.7), ry, Inches(2.8), rh, [{"t": a, "size": 15, "bold": True, "color": NAVY}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.5), ry, Inches(4.0), rh, [{"t": b, "size": 14, "color": GREY, "line_spacing": 1.05}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.6), ry, Inches(5.0), rh, [{"t": c, "size": 14, "color": DARK, "bold": True, "line_spacing": 1.05}],
        anchor=MSO_ANCHOR.MIDDLE)
    ry += rh + Inches(0.12)
rect(s, Inches(0.6), Inches(5.65), Inches(12.1), Inches(1.1), RGBColor(0xFC, 0xEF, 0xEF), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(5.65), Pt(6), Inches(1.1), RED)
txt(s, Inches(0.85), Inches(5.78), Inches(11.7), Inches(0.9),
    [{"t": "成果對比", "size": 13, "bold": True, "color": RED, "space_after": 3},
     {"t": [{"t": "SACF 在 MAE（0.5840）與 Corr（0.8691）優於 MGT（0.659 / 0.822）；Acc-7 達 ", "size": 14.5, "color": DARK},
            {"t": "53.06%，相對 MGT 提升 +2.62%", "size": 14.5, "bold": True, "color": RED}]}])
page(s)


# =====================================================================
# CHAPTER 3 — SACF 模型
# =====================================================================
divider("3", "SACF 情感感知跨模態融合模型", "Methodology",
        ["多分支單一模型 × 階層式跨模態融合", "多工損失 · 兩階段訓練", "零洩漏推斷流程"])

# 3.0 框架概覽
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第三章 · 框架概覽", "核心設計：把「集成多樣性」內建於單一模型")
txt(s, Inches(0.6), Inches(1.25), Inches(6.1), Inches(5.6),
    [bullet("共享上游編碼層，下游 4 條結構同形、參數獨立的並行分支", 15),
     bullet("各分支施以不同隨機失活率 0.10 / 0.20 / 0.30 / 0.40", 14, sub=True),
     bullet("走過不同隨機路徑 → 學到互補決策邊界", 14, sub=True),
     bullet("階層式跨模態融合（二階段）", 15, color=NAVY, bold=True),
     bullet("Stage 1：語言-音訊、語言-視覺初步對齊", 14, sub=True),
     bullet("Stage 2：低學習率引入跨模態對比損失深層對齊", 14, sub=True),
     bullet("三層方差降低推斷", 15, color=NAVY, bold=True),
     bullet("TTA×5 → 三次訓練參數空間平均 → 分類-回歸機率融合", 14, sub=True),
     bullet("嚴格零資料洩漏：融合超參數皆先驗設定", 14.5, color=RED, bold=True)])
# 右側成果指標
metric_chip(s, Inches(7.0), Inches(1.35), Inches(2.7), "53.06%", "Acc-7（主要指標）")
metric_chip(s, Inches(9.95), Inches(1.35), Inches(2.7), "0.5840", "MAE（越低越好）", vcolor=BLUE)
metric_chip(s, Inches(7.0), Inches(2.65), Inches(2.7), "0.8691", "Pearson Corr", vcolor=TEAL)
metric_chip(s, Inches(9.95), Inches(2.65), Inches(2.7), "85.42%", "Acc-2 二分類", vcolor=NAVY)
rect(s, Inches(7.0), Inches(4.0), Inches(5.65), Inches(2.5), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.2), Inches(4.12), Inches(5.3), Inches(0.4),
    [{"t": "推斷僅需一次前向傳播", "size": 14, "bold": True, "color": NAVY}])
txt(s, Inches(7.2), Inches(4.5), Inches(5.3), Inches(1.9),
    [{"t": "傳統 test-time 集成需維護多個權重檔、佔用大量顯存；", "size": 13, "color": DARK, "space_after": 4, "line_spacing": 1.1},
     {"t": "本研究將三次獨立訓練的權重於參數空間以 0.25 / 0.45 / 0.30 加權合併為單一檔案", "size": 13, "color": DARK, "space_after": 4, "line_spacing": 1.1},
     {"t": "4 分支於同一 forward 計算，開銷僅約單分支 1.4 倍", "size": 13, "bold": True, "color": BLUE, "line_spacing": 1.1}])
page(s)


# 共用：大圖頁
def figure_slide(chapter, title, img, caption, bullets, img_ratio=0.6):
    s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
    content_header(s, chapter, title)
    img_w = Inches(13.333 * img_ratio - 0.3)
    pic_fit(s, os.path.join(ASSETS, img), Inches(0.5), Inches(1.35),
            img_w, Inches(5.0))
    txt(s, Inches(0.5), Inches(6.42), img_w, Inches(0.4),
        [{"t": caption, "size": 12, "color": GREY, "bold": True}], align=PP_ALIGN.CENTER)
    bx = Inches(0.5) + img_w + Inches(0.25)
    bw = SW - bx - Inches(0.4)
    rect(s, bx, Inches(1.35), bw, Inches(5.45), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bx + Inches(0.2), Inches(1.55), bw - Inches(0.4), Inches(5.2),
        bullets)
    page(s)
    return s

# 圖1 整體架構
figure_slide("第三章 · 模型架構", "圖 1｜SACF 整體架構", "image3.png",
    "圖 1　SACF 整體架構（四橫向區塊）",
    [{"t": "四個橫向區塊", "size": 16, "bold": True, "color": NAVY, "space_after": 8},
     bullet("輸入：文字句子 + 語音音訊 + 人臉視覺序列", 14),
     bullet("共享編碼器：24 層 Transformer + 雙向 LSTM", 14),
     bullet("4 條並行分支（核心）", 14, color=RED, bold=True),
     {"t": "    詞元級顯著性閘控 → 二段式跨模態融合 → 共享投影 → 多工頭", "size": 12.5, "color": GREY, "space_after": 8, "line_spacing": 1.1},
     bullet("輸出：7 分類 + 2 分類 + 強度回歸", 14),
     bullet("分支平均 → TTA → 參數空間集成 → 機率融合", 14)],
    img_ratio=0.5)

# 資料集 + 圖2
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第三章 · 資料集", "CMU-MOSI 資料集與七類情感分布")
txt(s, Inches(0.6), Inches(1.3), Inches(5.9), Inches(5.4),
    [{"t": "CMU-MOSI 多模態情感分析標準基準", "size": 16, "bold": True, "color": NAVY, "space_after": 8},
     bullet("93 位 YouTube 評論者、2,199 個語句單位", 14.5),
     bullet("5 位標注員評分，平均落於 [−3, +3] 連續強度", 14.5),
     bullet("文字（轉錄）+ 音訊（COVAREP 5維）+ 視覺（FACET 20維）", 14.5),
     bullet("訓練+驗證合併 n=1,513；測試 686（僅評估一次）", 14.5, color=BLUE, bold=True),
     {"t": "分布偏移挑戰", "size": 15, "bold": True, "color": RED, "space_before": 8, "space_after": 6},
     bullet("測試集顯著偏向負面端（類 −3 約為訓練集 3 倍）", 14),
     bullet("→ 設計軟序數標籤、EMD 損失、三次集成多重正則緩解", 14, sub=True)])
pic_fit(s, os.path.join(ASSETS, "image4.png"), Inches(6.7), Inches(1.6), Inches(6.2), Inches(3.8))
txt(s, Inches(6.7), Inches(5.5), Inches(6.2), Inches(0.4),
    [{"t": "圖 2　CMU-MOSI 七類情感分布（訓練集 vs 測試集）", "size": 12, "color": GREY, "bold": True}],
    align=PP_ALIGN.CENTER)
page(s)

# 模型架構 — 共享編碼層 + 4分支
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第三章 · 模型架構", "共享編碼層　×　4 個並行分支的多樣性來源")
card(s, Inches(0.6), Inches(1.35), Inches(5.95), Inches(2.55), "共享編碼層（上游）",
     ["文字骨幹：microsoft/deberta-v3-large",
      "　24 層 Transformer · 隱藏維度 1,024 · 約 400M 參數",
      "音訊/視覺：各 2 層雙向 LSTM（前後向各 128 維）",
      "　投影至 d_modal = 128，僅對有效幀計算梯度",
      "共享而非每分支獨立 → 避免模型膨脹數倍"], headcolor=BLUE)
card(s, Inches(6.75), Inches(1.35), Inches(5.95), Inches(2.55), "4 個並行分支：三項多樣性機制",
     ["① 不同隨機失活率 0.10 / 0.20 / 0.30 / 0.40",
      "　等價於四個不同正則強度的子網路",
      "② 獨立的 PEA、跨模態融合與投影層參數",
      "③ 七分類頭初始化加入 0.005·(i+1) 高斯擾動",
      "　避免分支塌陷至近似解"], headcolor=TEAL)
rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.4), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.8), Inches(4.35), Inches(11.7), Inches(0.45),
    [{"t": "分支內部處理流程（單分支 forward）", "size": 15, "bold": True, "color": NAVY}])
steps = [("PEA\n極性增強注意力", BLUE), ("Top-K 情感詞元\n選擇 K=5", TEAL),
         ("階層式跨模態\n融合（二階段）", NAVY), ("共享投影層\n1024→512", GOLD),
         ("多工頭\n7分類/2分類/回歸", RED)]
sx = Inches(0.85); sw = Inches(2.05); sy = Inches(4.95); sh2 = Inches(1.15); g = Inches(0.28)
for i, (t, c) in enumerate(steps):
    x = sx + i * (sw + g)
    rect(s, x, sy, sw, sh2, c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, sy, sw, sh2, [{"t": t, "size": 13, "bold": True, "color": WHITE, "line_spacing": 1.05}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 4:
        txt(s, x + sw, sy, g, sh2, [{"t": "▶", "size": 14, "color": GREY}],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page(s)

# 圖3 PEA
figure_slide("第三章 · 核心模組", "圖 3｜極性增強注意力 PEA", "image5.png",
    "圖 3　極性增強注意力（PEA）模組",
    [{"t": "PEA：為每個詞元學習情感顯著性", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     bullet("對每個 DeBERTa 詞元學 0~1 的情感顯著性閘值 gᵢ", 13.5),
     bullet("以閘值加權平均池化，產出集中於情感詞的句向量", 13.5),
     bullet("保守係數 0.75/0.25：閘值全 0 時仍保留 75% 原始訊號", 13.5),
     bullet("避免閘值梯度初期不穩定破壞語境", 13, sub=True),
     bullet("輸出 xₗ 供共享投影層；閘值序列 g 供下游 Top-K 選詞", 13.5)],
    img_ratio=0.58)

# 圖4 跨模態融合四步驟
figure_slide("第三章 · 核心模組", "圖 4｜情感感知跨模態融合四步驟", "image6.png",
    "圖 4　情感感知跨模態融合四步驟",
    [{"t": "以「情感感知查詢」取代首詞元查詢", "size": 14.5, "bold": True, "color": NAVY, "space_after": 8},
     bullet("步驟1：依顯著性取 Top-K=5 詞元，濾除停用詞", 13.5),
     bullet("步驟2：注意力加權彙整為情感感知查詢 q_sa", 13.5),
     bullet("步驟3：音訊/視覺投影至語言空間作鍵值，縮放點積注意力提取跨模態證據", 13.5),
     bullet("步驟4：FFN + sigmoid 閘控 + LayerNorm 殘差融合", 13.5),
     {"t": "→ 4 分支融合模組參數完全獨立，注意力分布不同，是內部集成多樣性的關鍵來源", "size": 12.5, "color": RED, "bold": True, "space_before": 6, "line_spacing": 1.1}],
    img_ratio=0.58)

# 圖5 四分支多樣性與內部集成
figure_slide("第三章 · 內部集成", "圖 5｜四條並行分支多樣性與內部集成", "image7.png",
    "圖 5　四條並行分支之多樣性來源與內部集成",
    [{"t": "三層方差降低的第一層", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     bullet("各分支單獨 Acc-7：49.56% ~ 50.29%", 14),
     bullet("4 分支內部平均後 ≈ 50.00%", 14),
     bullet("再經參數空間平均 + 機率融合", 14, sub=True),
     {"t": [{"t": "→ 最終 Acc-7 提升至 ", "size": 14, "color": DARK},
            {"t": "53.06%", "size": 16, "bold": True, "color": RED}], "space_before": 6, "space_after": 8},
     bullet("計算開銷僅單分支約 1.4 倍，大幅降低 GPU 成本", 13.5)],
    img_ratio=0.58)

# 圖6 多工損失
figure_slide("第三章 · 損失函數", "圖 6｜多工損失組成結構", "image8.png",
    "圖 6　多工損失組成結構（三層）",
    [{"t": "內部驅動的多工損失（三層）", "size": 14.5, "bold": True, "color": NAVY, "space_after": 7},
     bullet("Layer 1 單分支：軟序數 CE + 序數 EMD + 二元 CE + Smooth-L1", 13),
     bullet("Layer 2 跨分支：聚合 + 特徵多樣性 + 跨模態對比（僅 Stage 2）", 13),
     bullet("Layer 3：R-Drop 一致性正則 + 總損失 L_total", 13),
     {"t": "關鍵設計", "size": 13.5, "bold": True, "color": RED, "space_before": 7, "space_after": 4},
     bullet("軟序數標籤：相鄰類錯誤懲罰 < 遠距類", 12.5, sub=True),
     bullet("EMD：以 CDF 差異約束分布形狀", 12.5, sub=True),
     bullet("InfoNCE 對比：拉近同樣本三模態嵌入（τ=0.07）", 12.5, sub=True)],
    img_ratio=0.55)

# 圖7 兩階段訓練
figure_slide("第三章 · 訓練策略", "圖 7｜兩階段訓練協議全景", "image9.png",
    "圖 7　兩階段訓練協議全景",
    [{"t": "Stage 1 基底訓練（60 epoch）", "size": 14, "bold": True, "color": BLUE, "space_after": 5},
     bullet("Phase1 凍結 DeBERTa 下層6層（E1–20）", 12.5, sub=True),
     bullet("Phase2 全模型微調（E20–42）", 12.5, sub=True),
     bullet("Phase3 SWA 隨機權重平均（E42–60）", 12.5, sub=True),
     {"t": "Stage 2 CMC 對比精修（20 epoch）", "size": 14, "bold": True, "color": TEAL, "space_before": 6, "space_after": 5},
     bullet("低學習率 + 加入跨模態對比損失", 12.5, sub=True),
     {"t": "三次獨立執行（seed/協議不同）", "size": 13.5, "bold": True, "color": RED, "space_before": 6, "space_after": 4},
     bullet("參數空間加權 0.25 / 0.45 / 0.30 合併為單一檔", 12.5, sub=True)],
    img_ratio=0.6)

# 圖3-8 零洩漏推斷
figure_slide("第三章 · 推斷流程", "圖 3-8｜零洩漏推斷流程", "image10.png",
    "圖 3-8　零洩漏推斷流程（三層方差降低）",
    [{"t": "三層方差降低設計", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     bullet("① TTA×5：保留 dropout 隨機性，5 次前向取平均", 14),
     bullet("② 參數空間平均：三次訓練已預先合併為單檔", 14),
     bullet("③ 分類-回歸機率融合：對數空間幾何平均", 14),
     {"t": "所有融合超參數（α=0.65, σ=0.65, T=1.0）皆先驗設定，不依測試集調整 → 嚴格零資料洩漏",
      "size": 13, "color": RED, "bold": True, "space_before": 8, "line_spacing": 1.15}],
    img_ratio=0.58)

# 圖3-9 機率融合
figure_slide("第三章 · 推斷流程", "圖 3-9｜分類-回歸機率融合示意", "image11.png",
    "圖 3-9　分類-回歸機率融合示意（樣本索引 316）",
    [{"t": "互補資訊融合", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     bullet("(a) 分類頭輸出的七類機率分布", 14),
     bullet("(b) 回歸頭預測經高斯核轉換的七類分布", 14),
     bullet("(c) 對數空間幾何平均合併的最終分布", 14),
     {"t": [{"t": "融合後 Acc-7 提升 ", "size": 14, "color": DARK},
            {"t": "+0.44%（52.62% → 53.06%）", "size": 15, "bold": True, "color": RED}],
      "space_before": 8}],
    img_ratio=0.62)


# =====================================================================
# CHAPTER 4 — 實驗結果
# =====================================================================
divider("4", "SACF 情感評分模型實驗結果", "Experimental Results",
        ["CMU-MOSI 五指標全面比較", "多種子穩健性 · 集成效益", "跨資料集泛化與參數效率"])

# 評估指標
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第四章 · 評估指標", "五項評估指標說明")
mets = [("Acc-7", "七分類準確率", "連續值四捨五入至 −3~+3\n細粒度識別主要指標", RED),
        ("Acc-2", "二分類準確率", "以正負為界判斷極性\n反映基本極性辨別能力", BLUE),
        ("F1", "加權 F1 分數", "對類別不平衡適應性強\n與 Acc-2 共評極性品質", TEAL),
        ("MAE", "平均絕對誤差", "預測與真實的平均絕對差\n越低越好（回歸精度）", GOLD),
        ("Corr", "Pearson 相關係數", "預測與真實的線性相關\n[−1,1]，越高越好", NAVY)]
mx = Inches(0.6); mw = Inches(2.3); my = Inches(1.7); mh = Inches(2.6); mg = Inches(0.15)
for i, (k, n, d, c) in enumerate(mets):
    x = mx + i * (mw + mg)
    rect(s, x, my, mw, mh, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, my, mw, Inches(0.7), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, my + Inches(0.04), mw, Inches(0.62), [{"t": k, "size": 21, "bold": True, "color": WHITE, "name": FONT_NUM}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, my + Inches(0.85), mw, Inches(0.45), [{"t": n, "size": 14, "bold": True, "color": DARK}],
        align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.12), my + Inches(1.35), mw - Inches(0.24), Inches(1.1),
        [{"t": d, "size": 12, "color": GREY, "line_spacing": 1.15}], align=PP_ALIGN.CENTER)
txt(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.8),
    [{"t": "基線比較方法（參考 MGT 整理的 CMU-MOSI 對照表，共 13 個代表性方法）", "size": 14.5, "bold": True, "color": NAVY, "space_after": 8},
     {"t": "LF-DNN · Graph-MFN · MFM · MISA · Self-MM · MMIM · DMD · ITHP · MCL-MCF · ConFEDE · CLGSI · UniMSE · MGT",
      "size": 14, "color": DARK, "line_spacing": 1.2, "space_after": 6},
     {"t": "其中 MGT（IEEE TAFFC 2025）以 PLM 為主模態 + 跨模態加法注意力，是本研究主要對標的最強基線。",
      "size": 13.5, "color": RED}])
page(s)

# CMU-MOSI 比較 + 圖8
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第四章 · 主要結果", "CMU-MOSI 效能比較：四項指標達 SOTA")
pic_fit(s, os.path.join(ASSETS, "image12.png"), Inches(0.5), Inches(1.35), Inches(7.0), Inches(4.9))
txt(s, Inches(0.5), Inches(6.4), Inches(7.0), Inches(0.4),
    [{"t": "圖 8　CMU-MOSI 多指標比較（前八模型，紅色為本研究）", "size": 11.5, "color": GREY, "bold": True}],
    align=PP_ALIGN.CENTER)
# 右側對比表
bx = Inches(7.75); bw = Inches(5.0)
rows = [("指標", "MGT 基線", "SACF", True),
        ("Acc-7 ↑", "50.44%", "53.06%", False),
        ("Acc-2 ↑", "86.30%", "85.42%", False),
        ("F1 ↑", "86.28%", "85.41%", False),
        ("MAE ↓", "0.659", "0.5840", False),
        ("Corr ↑", "0.822", "0.8691", False)]
ry = Inches(1.45); rh = Inches(0.62)
for i, (a, b, c, hd) in enumerate(rows):
    y = ry + i * (rh + Inches(0.07))
    bg = NAVY if hd else (WHITE if i % 2 else CARDBG)
    rect(s, bx, y, bw, rh, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tc = WHITE if hd else DARK
    win = c in ("53.06%", "0.5840", "0.8691")  # SACF 勝出
    txt(s, bx + Inches(0.15), y, Inches(1.7), rh, [{"t": a, "size": 14, "bold": True, "color": tc}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, bx + Inches(1.85), y, Inches(1.45), rh, [{"t": b, "size": 14, "color": tc, "name": FONT_NUM}], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    txt(s, bx + Inches(3.3), y, Inches(1.6), rh, [{"t": c, "size": 14.5, "bold": True, "color": (WHITE if hd else (RED if win else DARK)), "name": FONT_NUM}], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
txt(s, bx, Inches(5.85), bw, Inches(1.0),
    [bullet("Acc-7 +2.62%、MAE −11.4%、Corr +0.047 全面領先", 13, color=RED, bold=True, sa=4),
     bullet("Acc-2/F1 僅微幅落後 0.88%（多模態互補訊號）", 12.5)])
page(s)

# 多種子穩健性 + 集成效益
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第四章 · 穩健性與集成", "多種子穩健性分析 與 集成效益")
card(s, Inches(0.6), Inches(1.35), Inches(5.95), Inches(2.9), "三次獨立執行（單模型 Acc-7）",
     ["Run 1（seed=42，標準協議）：52.62%",
      "Run 2（seed=42，Stage 2 延長）：52.48%",
      "Run 3（seed=5678，加兩輪 BAN）：51.60%",
      "標準差 ≈ 0.43% → 高度一致",
      "三次皆超越 MGT 基線約 1–2%",
      "MAE 0.591–0.609、Corr 均 > 0.860"], headcolor=BLUE)
card(s, Inches(6.75), Inches(1.35), Inches(5.95), Inches(2.9), "三種子等權集成效益",
     ["每樣本 3 種子 × 5 TTA = 15 次獨立推斷",
      "Acc-7：53.06%（較最優單執行 +0.44%）",
      "　　　（較最弱單執行 +1.46%）",
      "MAE：0.5840（較單執行均值降 ≈1.6%）",
      "Corr：0.8691（高於所有單執行）",
      "不同種子探索不同收斂路徑 → 互補決策邊界"], headcolor=TEAL)
rect(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.1), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.8), Inches(4.68), Inches(11.7), Inches(0.45),
    [{"t": "三層方差降低如何疊加出最終 53.06%", "size": 15, "bold": True, "color": NAVY}])
chain = [("單分支\n~50.0%", GREY), ("內部集成\n50.00%", BLUE),
         ("參數空間平均\n+ 機率融合", TEAL), ("最終 Acc-7\n53.06%", RED)]
cxx = Inches(1.2); cww = Inches(2.5); cyy = Inches(5.3); chh = Inches(1.0); cgg = Inches(0.55)
for i, (t, c) in enumerate(chain):
    x = cxx + i * (cww + cgg)
    rect(s, x, cyy, cww, chh, c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, cyy, cww, chh, [{"t": t, "size": 14, "bold": True, "color": WHITE, "line_spacing": 1.05}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        txt(s, x + cww, cyy, cgg, chh, [{"t": "▶", "size": 18, "color": GREY}],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page(s)

# 跨資料集 MMAFFBen + 圖9
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第四章 · 泛化性", "跨資料集泛化（一）：MMAFFBen 五任務")
pic_fit(s, os.path.join(ASSETS, "image13.png"), Inches(0.5), Inches(1.4), Inches(7.0), Inches(4.4))
txt(s, Inches(0.5), Inches(5.95), Inches(7.0), Inches(0.4),
    [{"t": "圖 9　SACF-Text 與其他語言模型在 MMAFFBen 五任務平均 macro-F1", "size": 11, "color": GREY, "bold": True}],
    align=PP_ALIGN.CENTER)
txt(s, Inches(7.75), Inches(1.4), Inches(5.0), Inches(5.4),
    [{"t": "輕量遷移學習設定", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     bullet("以 CMU-MOSI 最佳檢查點初始化，僅微調 1–3 epoch", 14),
     bullet("平均 ma-F1 = 56.74，排名第三", 14.5, color=RED, bold=True),
     bullet("僅次於 MMAFFLM-3b（58.2）/ 7b（57.1）", 13.5, sub=True),
     bullet("以 0.435B 參數達同級水準", 13.5, sub=True),
     bullet("MMS 三語情感強度 61.42 → 排名第一", 14),
     bullet("Onlineshopping 中文極性 92.82，領先 GPT-4o-mini +30.9%", 14),
     bullet("大幅領先 GPT-4o-mini（51.8）與 EmoLlama-7b（37.8）", 13.5)])
page(s)

# CMU-MOSEI 圖10 + 圖11
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第四章 · 泛化性", "跨資料集泛化（二）：CMU-MOSEI（10× 規模）")
pic_fit(s, os.path.join(ASSETS, "image14.png"), Inches(0.45), Inches(1.35), Inches(6.2), Inches(3.1))
txt(s, Inches(0.45), Inches(4.5), Inches(6.2), Inches(0.35),
    [{"t": "圖 10　CMU-MOSEI 七分類準確率模型排序", "size": 10.5, "color": GREY, "bold": True}], align=PP_ALIGN.CENTER)
pic_fit(s, os.path.join(ASSETS, "image15.png"), Inches(0.45), Inches(4.95), Inches(6.2), Inches(2.0))
txt(s, Inches(0.45), Inches(7.0), Inches(6.2), Inches(0.3),
    [{"t": "圖 11　CMU-MOSEI 多指標比較", "size": 10.5, "color": GREY, "bold": True}], align=PP_ALIGN.CENTER)
txt(s, Inches(6.95), Inches(1.35), Inches(5.85), Inches(5.5),
    [{"t": "僅用文字模態，超越多模態基線", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     bullet("22,852 語句，規模為 CMU-MOSI 十倍以上", 14),
     bullet("Acc-7 = 55.49% → 13 個對照模型中排名第一", 14.5, color=RED, bold=True),
     bullet("略高於 MGT（55.10%），優於 DMD/UniMSE/MMIM", 13.5, sub=True),
     bullet("MAE 0.4881、Corr 0.8243 皆最佳", 14),
     bullet("Within-1 達 96.24% → 序數擬合品質極佳", 14),
     {"t": "文字分支已內化 CMU-MOSI 習得的情感先驗，能有效遷移至更大規模資料集；驗證 PEA + 共享編碼 + 多工損失將情感歸納偏好內化於語言骨幹。",
      "size": 13, "color": DARK, "line_spacing": 1.15, "space_before": 8}])
page(s)

# SemEval + 參數效率
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第四章 · 泛化性與效率", "SemEval-2018 三語言 與 參數效率分析")
card(s, Inches(0.6), Inches(1.35), Inches(5.95), Inches(2.9), "SemEval-2018 Task 1（EN/AR/ES）",
     ["三語言整體分數 66.99，排名第三",
      "英語 73.04 → 超越 18× 大的 GPT-4o-mini（70.00）",
      "西班牙語 69.28，接近 EmoLlama-7b（76.37）",
      "阿拉伯語 58.64 偏弱（英文預訓練背景所致）",
      "驗證情感表示具跨語言遷移能力"], headcolor=BLUE)
card(s, Inches(6.75), Inches(1.35), Inches(5.95), Inches(2.9), "參數效率（分/十億參數）",
     ["SemEval：154.0 分/B → 全部模型最高",
      "　為第二名 MMAFFLM-3b（21.8）的 7.1 倍",
      "MMAFFBen：127.6 分/B",
      "　為 MMAFFLM-3b 的 6.6 倍、GPT-4o-mini 的 20.3 倍",
      "以 0.435B 參數媲美 16–26× 規模的模型"], headcolor=RED)
rect(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.1), RGBColor(0xEC, 0xF4, 0xEC), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(4.55), Pt(6), Inches(2.1), TEAL)
txt(s, Inches(0.85), Inches(4.7), Inches(11.7), Inches(1.9),
    [{"t": "核心主張獲得驗證", "size": 15, "bold": True, "color": TEAL, "space_after": 6},
     {"t": "在情感分析任務上，針對性設計的跨模態注意力架構（SACF + PEA）結合情感標記語料精調，能以極低參數代價實現媲美大型語言模型的跨任務情感泛化能力。",
      "size": 15, "color": DARK, "line_spacing": 1.25, "space_after": 6},
     {"t": "→ 對計算資源受限的部署場景（如本研究的本地多代理情感感知系統）具重要實際意義。",
      "size": 14, "bold": True, "color": NAVY}])
page(s)


# =====================================================================
# CHAPTER 5 — 應用
# =====================================================================
divider("5", "情緒感知多代理模擬系統", "Application System",
        ["感知-思考-行動-情緒 認知循環", "三維記憶 × 七類情緒子系統", "系統畫面實際展示"])

# 系統整體架構
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第五章 · 系統架構", "認知循環：在生成式代理上延伸情緒感知層")
txt(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.55),
    [{"t": [{"t": "承襲 Park 等人生成式代理（記憶流 / 反思 / 規劃），", "size": 15, "color": DARK},
            {"t": "新增「情緒（Emotion）」階段，並在地化為中文（Qwen2.5:32b 本地推理）", "size": 15, "bold": True, "color": NAVY}]}])
cyc = [("感知\nPercept", "視野範圍掃描\n重要性閾值過濾", BLUE),
       ("思考\nThink", "記憶檢索\n反思生成", TEAL),
       ("規劃\nPlan", "結合特徵目標\n生成行動序列", NAVY),
       ("行動\nAct", "環境互動\n對話輸出", GOLD),
       ("情緒\nEmotion", "情緒表徵/衰減\n行為調節", RED)]
cx = Inches(0.7); cw = Inches(2.3); cy = Inches(2.0); ch = Inches(1.5); g = Inches(0.18)
for i, (t, d, c) in enumerate(cyc):
    x = cx + i * (cw + g)
    rect(s, x, cy, cw, ch, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, cy, cw, Inches(0.62), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, cy, cw, Inches(0.62), [{"t": t, "size": 14, "bold": True, "color": WHITE, "line_spacing": 0.95}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, cy + Inches(0.72), cw, Inches(0.75), [{"t": d, "size": 12, "color": DARK, "line_spacing": 1.05}],
        align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, x + cw, cy, g, ch, [{"t": "▶", "size": 13, "color": GREY}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.6), Inches(3.85), Inches(12.1), Inches(3.0),
    [{"t": "相較原始框架的四項主要延伸", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     bullet("將 SACF 跨模態情感分析模型嵌入感知模組 → 感知他人言談、表情、語氣中的情感", 14.5),
     bullet("認知循環新增情緒階段 → 內在情緒狀態的表徵、衰減與行為調節", 14.5),
     bullet("調整提示工程模板 → 調節代理行動與行為", 14.5),
     bullet("在地化為中文語境 → 採本地模型 Qwen2.5:32b 作為核心推理引擎", 14.5),
     {"t": "五大組件：智能體核心引擎 · 記憶與情感管理 · 環境建模 · 語言模型介面 · 視覺化分析工具",
      "size": 13.5, "color": GREY, "space_before": 8, "bold": True}])
page(s)

# 記憶 + 情緒子系統
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第五章 · 核心子系統", "三維記憶管理　×　七類情緒狀態管理")
card(s, Inches(0.6), Inches(1.35), Inches(5.95), Inches(5.3), "記憶管理系統（仿人類多層記憶）",
     ["三類記憶：觀察記憶 · 反思記憶 · 對話記憶",
      "",
      "三維檢索評分（加權組合，權重 α 皆=1）：",
      "　• 重要性：LLM 對事件 1–10 分評估",
      "　　（刷牙低分 / 分手升學高分）",
      "　• 時效性：指數衰減，衰減因子 0.995",
      "　• 相關性：查詢與記憶的嵌入餘弦相似度",
      "",
      "動態更新：訪問即更新時間戳",
      "定期清理：預設記憶保存 30 天",
      "→ 既優先近期情感記憶，又不喪失重要事件追蹤"], headcolor=BLUE)
card(s, Inches(6.75), Inches(1.35), Inches(5.95), Inches(5.3), "情緒子系統（表徵→感知→生成→衰減→作用）",
     ["七類離散情緒：憤怒/厭惡/恐懼/喜悅/中性/悲傷/驚訝",
      "　強度 0~1 連續尺度，記錄原因與持續時間",
      "",
      "跨模態情感感知：SACF 整合語言/語音/視覺",
      "　降級路徑：僅文字 → LLM 文字推測",
      "　強度—類別投影：對應七類離散情緒",
      "",
      "生成更新：候選情緒需達觸發閾值才轉換（慣性）",
      "自然衰減：固定速率遞減，低於門檻回歸中性",
      "作用範圍：記憶形成/檢索、對話語氣、規劃偏好"], headcolor=RED)
page(s)

# 系統畫面展示 1
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第五章 · 系統展示", "畫面展示（一）：社交網絡與角色資料")
pic_fit(s, os.path.join(ASSETS, "image16.png"), Inches(0.5), Inches(1.45), Inches(5.9), Inches(4.3))
txt(s, Inches(0.5), Inches(5.85), Inches(5.9), Inches(0.4),
    [{"t": "居民關係力導向圖（節點大小=活躍度，連線粗細=互動頻率）", "size": 11.5, "color": GREY, "bold": True}], align=PP_ALIGN.CENTER)
pic_fit(s, os.path.join(ASSETS, "image17.png"), Inches(6.7), Inches(1.45), Inches(6.1), Inches(4.3))
txt(s, Inches(6.7), Inches(5.85), Inches(6.1), Inches(0.4),
    [{"t": "角色詳細資料：人格特質 · 情緒狀態 · 社交關係統計", "size": 11.5, "color": GREY, "bold": True}], align=PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
    [{"t": "力導向圖 + 對話長度矩陣熱力圖 + 行為時間軸 + 情感狀態變化圖，為湧現行為提供量化分析。",
      "size": 13, "color": NAVY, "bold": True}], align=PP_ALIGN.CENTER)
page(s)

# 系統畫面展示 2
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第五章 · 系統展示", "畫面展示（二）：智能問卷調查系統")
pic_fit(s, os.path.join(ASSETS, "image18.png"), Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.7), border=True)
txt(s, Inches(0.5), Inches(3.12), Inches(12.3), Inches(0.32),
    [{"t": "問卷管理系統：支援建立 / 複製 / 批量操作 / 版本管理 / AI 自動填寫", "size": 11.5, "color": GREY, "bold": True}], align=PP_ALIGN.CENTER)
pic_fit(s, os.path.join(ASSETS, "image19.png"), Inches(0.5), Inches(3.5), Inches(8.4), Inches(3.4))
txt(s, Inches(0.5), Inches(6.95), Inches(8.4), Inches(0.3),
    [{"t": "社區居民情緒感知問卷：基於情緒記憶與互動歷史自動填答", "size": 11, "color": GREY, "bold": True}], align=PP_ALIGN.CENTER)
txt(s, Inches(9.1), Inches(3.6), Inches(3.7), Inches(3.3),
    [{"t": "AI 智能問卷新範式", "size": 14.5, "bold": True, "color": NAVY, "space_after": 8},
     bullet("單選/多選/評分/開放文字", 13),
     bullet("整合情緒記憶檢索自動填答", 13),
     bullet("回應一致且可追溯", 13),
     bullet("可反覆採樣支援假設驗證", 13),
     bullet("匯出 CSV / JSON / Excel", 13)])
page(s)


# =====================================================================
# CHAPTER 6 — 結論
# =====================================================================
divider("6", "結論與貢獻", "Conclusion & Contributions",
        ["研究成果總結", "三大主要貢獻", "研究限制與未來展望"])

# 主要貢獻
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第六章 · 主要貢獻", "三大主要貢獻")
contrib = [
    ("貢獻一", "情感顯著詞元驅動的跨模態融合新範式",
     "以 sentiment-salient tokens 構建情感感知查詢，使跨模態注意力聚焦於與情感極性最相關的語意單元。", BLUE),
    ("貢獻二", "適用序數情感標籤的多工損失與零洩漏推斷",
     "序數 EMD + 軟序數標籤 + R-Drop + InfoNCE 兼顧分類與回歸；融合超參數先驗設定，嚴格避免測試時資訊洩漏。", TEAL),
    ("貢獻三", "系統性驗證跨語言、跨任務的廣泛適用性",
     "於 MMAFFBen、SemEval-2018、CMU-MOSEI 驗證強泛化與優異參數效率，並落地為情緒感知多代理模擬系統。", RED),
]
cy = Inches(1.4)
for no, t, d, c in contrib:
    rect(s, Inches(0.6), cy, Inches(12.1), Inches(1.45), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.6), cy, Inches(2.0), Inches(1.45), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.6), cy, Inches(2.0), Inches(1.45), [{"t": no, "size": 22, "bold": True, "color": WHITE}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.85), cy + Inches(0.18), Inches(9.6), Inches(1.1),
        [{"t": t, "size": 17, "bold": True, "color": DARK, "space_after": 5},
         {"t": d, "size": 13.5, "color": GREY, "line_spacing": 1.15}])
    cy += Inches(1.62)
page(s)

# 成果總結 + 限制與未來
s = add_slide(); rect(s, 0, 0, SW, SH, LIGHT)
content_header(s, "第六章 · 總結與展望", "研究成果總結　與　限制及未來展望")
# 成果指標列
metric_chip(s, Inches(0.6), Inches(1.35), Inches(2.85), "53.06%", "CMU-MOSI Acc-7", RED)
metric_chip(s, Inches(3.65), Inches(1.35), Inches(2.85), "55.49%", "CMU-MOSEI Acc-7", BLUE)
metric_chip(s, Inches(6.7), Inches(1.35), Inches(2.85), "154 分/B", "SemEval 參數效率", TEAL)
metric_chip(s, Inches(9.75), Inches(1.35), Inches(2.95), "9 位", "情緒感知 AI 居民", GOLD)
card(s, Inches(0.6), Inches(2.75), Inches(5.95), Inches(3.9), "研究限制",
     ["硬體資源：單一 LLM 並行，居民限 9 位以內",
      "語言模型能力邊界：複合/矛盾情緒仍受限",
      "封閉系統：9 個固定角色，外部效度需謹慎",
      "觀察時程有限：長期情感演化難以充分驗證",
      "Acc-7 53.06% 仍有空間（標記主觀不一致的天花板效應）"], headcolor=GREY, bg=RGBColor(0xF0, 0xF0, 0xF2))
card(s, Inches(6.75), Inches(2.75), Inches(5.95), Inches(3.9), "未來展望",
     ["更大規模預訓練語言模型與跨模態預訓練資料",
      "多語多標籤骨幹改善 AR / XED 等弱項任務",
      "擴大代理數量與更長觀察週期",
      "教育輔助、心理健康、市場研究、社會政策評估",
      "情感感知計算社會科學的新方法論框架"], headcolor=TEAL)
page(s)

# 致謝 / Q&A
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.7), SW, Pt(4), GOLD)
rect(s, 0, Inches(4.7), SW, Pt(4), GOLD)
txt(s, Inches(0.6), Inches(2.95), Inches(12.1), Inches(1.6),
    [{"t": "敬請指教", "size": 52, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 10},
     {"t": "Thank You　|　Q & A", "size": 22, "color": GOLD, "align": PP_ALIGN.CENTER}])
txt(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.2),
    [{"t": "研究生：李昇峰　　指導教授：高君豪 教授", "size": 16, "color": RGBColor(0xD9, 0xE2, 0xF0), "align": PP_ALIGN.CENTER, "space_after": 6},
     {"t": "淡江大學 統計與資料科學碩士班　·　中華民國 115 年 6 月", "size": 13, "color": RGBColor(0xAE, 0xBC, 0xD4), "align": PP_ALIGN.CENTER}])
page(s)


out = os.path.join(HERE, "李昇峰_碩士論文口試簡報.pptx")
prs.save(out)
print("已生成：", out, "  共", len(prs.slides.__iter__.__self__._sldIdLst), "頁")
